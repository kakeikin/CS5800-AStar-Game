"""
Build A* poster as editable PPTX (36×24 inches).
All text, colors, and shapes are real PPTX objects — fully editable in PowerPoint/Keynote.

Run: python3 build_poster_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Scale: poster is 1440×960 px → 36×24 inches ─────────────────────────────
S = 36 / 1440  # inches per pixel

def px(n):      return Inches(n * S)
def fsz(n):     return Pt(round(n * 1.35))   # px → pt at poster print scale

# ── Colors ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1a, 0x23, 0x7e)
LAV     = RGBColor(0xe8, 0xea, 0xf6)
LAV_MID = RGBColor(0xc5, 0xca, 0xe9)
PAGE    = RGBColor(0xdd, 0xe1, 0xf0)
CARD    = RGBColor(0xff, 0xff, 0xff)
TEXT    = RGBColor(0x1c, 0x1c, 0x2e)
MUTED   = RGBColor(0x4a, 0x4a, 0x6a)
DIJK    = RGBColor(0x02, 0x84, 0xc7)
MANH    = RGBColor(0x05, 0x96, 0x69)
EUCL    = RGBColor(0xdc, 0x68, 0x03)
AMBER   = RGBColor(0xd9, 0x77, 0x06)
AMBERL  = RGBColor(0xff, 0xf3, 0xcd)
PATH_C  = RGBColor(0x7c, 0x3a, 0xed)
WHITE   = RGBColor(0xff, 0xff, 0xff)

# ── Poster shell ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = px(1440)
prs.slide_height = px(960)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def rect(l, t, w, h, fill=None, line_color=None, line_w_pt=0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w_pt)
    else:
        shp.line.fill.background()
    return shp

def tb(text, l, t, w, h, size, bold=False, italic=False,
       color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = fsz(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def tb_multi(lines, l, t, w, h, size, color=MUTED, spacing_pt=6):
    """Multi-line text box with multiple paragraphs."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing_pt) if i > 0 else Pt(0)
        run = p.add_run()
        run.text = txt
        run.font.size  = fsz(size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return box

def sec_card(l, t, w, h):
    """White card with lavender header strip."""
    rect(l, t, w, h, fill=CARD, line_color=LAV_MID, line_w_pt=0.5)

def sec_hdr(text, l, t, w):
    """Section header strip."""
    hdr_h = px(29)
    rect(l, t, w, hdr_h, fill=LAV)
    tb(text.upper(), l + px(8), t + px(4), w - px(16), hdr_h - px(6),
       size=9, bold=True, color=NAVY)
    return t + hdr_h  # returns y after header

# ═════════════════════════════════════════════════════════════════════════════
# BACKGROUND
# ═════════════════════════════════════════════════════════════════════════════
rect(0, 0, px(1440), px(960), fill=PAGE)

# ═════════════════════════════════════════════════════════════════════════════
# HEADER  (0, 0) → (1440, ~138)
# ═════════════════════════════════════════════════════════════════════════════
HDR_H = px(138)
rect(0, 0, px(1440), HDR_H, fill=NAVY)

# Title
tb("A* Pathfinding Visualizer",
   px(30), px(18), px(700), px(65),
   size=44, bold=True, color=WHITE)

# Authors + affiliation
tb("Jiaxin Jia    ·    Xiaoyuan Lu    ·    Xinyuan Fan",
   px(30), px(88), px(700), px(24),
   size=15.5, color=RGBColor(0xc5, 0xca, 0xe9))
tb("Khoury College of Computer Sciences  ·  Northeastern University",
   px(30), px(113), px(700), px(22),
   size=13.5, italic=True, color=RGBColor(0x9f, 0xa8, 0xda))

# NEU logo text approximation (white on navy)
tb("Northeastern\nUniversity",
   px(1230), px(38), px(180), px(65),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# BODY GRID  top: 138+13=151, bottom: 960-67-13=880  → height 729
# 3 columns each ~462px wide, gap 13px, left padding 13px
# ═════════════════════════════════════════════════════════════════════════════
BODY_T = px(151)
COL_W  = px(462)
GAP    = px(13)
C1     = px(13)
C2     = C1 + COL_W + GAP
C3     = C2 + COL_W + GAP

# ─────────────────────────────────────────────────────────────────────────────
# LEFT COLUMN
# ─────────────────────────────────────────────────────────────────────────────

# — Abstract (flex:0 0 auto → ~146px) —
SEC_T = BODY_T
SEC_H = px(146)
sec_card(C1, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Abstract", C1, SEC_T, COL_W)
tb("An interactive browser-based tool running Dijkstra, A* Manhattan, and A* Euclidean "
   "side-by-side on weighted grids. Users paint walls (∞), grass (×2), and swamps (×5), "
   "then watch all three execute with per-cell f(n)/g(n)/h(n) overlays. Benchmarks "
   "confirm A* cuts node expansions by ~24% vs Dijkstra on identical paths.",
   C1 + px(10), body_t + px(6), COL_W - px(20), px(95),
   size=12, color=TEXT)

# — Background & Motivation (flex:0 0 auto → ~212px) —
SEC_T += SEC_H + GAP
SEC_H = px(212)
sec_card(C1, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Background & Motivation", C1, SEC_T, COL_W)
tb("Algorithm trade-offs are hard to internalize from pseudocode alone — the contrast "
   "between uninformed and informed search only becomes intuitive when you can watch "
   "both exploration frontiers side by side in real time.",
   C1 + px(10), body_t + px(6), COL_W - px(20), px(54),
   size=12, color=TEXT)
tb("Most tools show one algorithm at a time on uniform-cost grids, making direct "
   "comparison impossible. This project closes that gap: Dijkstra, A* Manhattan, and "
   "A* Euclidean run simultaneously on shared weighted grids, with step-by-step replay "
   "and per-cell f/g/h overlays that reveal exactly why each algorithm expands each node.",
   C1 + px(10), body_t + px(66), COL_W - px(20), px(72),
   size=12, color=TEXT)
tb("Student feedback during development confirmed the core gap: watching A*'s frontier "
   "stay focused while Dijkstra fans out in all directions builds intuition no static "
   "textbook diagram can replicate.",
   C1 + px(10), body_t + px(144), COL_W - px(20), px(54),
   size=12, color=TEXT)

# — Algorithm Overview (flex:0 0 auto → ~156px) —
SEC_T += SEC_H + GAP
SEC_H = px(156)
sec_card(C1, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Algorithm Overview", C1, SEC_T, COL_W)

# Table header
rect(C1, body_t, COL_W, px(20), fill=LAV)
col_ws = [px(130), px(145), px(175)]
col_xs = [C1 + px(6), C1 + px(138), C1 + px(290)]
for hdr_txt, cx_col, cw in zip(["Algorithm", "Heuristic h(n)", "Best suited for"],
                                col_xs, col_ws):
    tb(hdr_txt, cx_col, body_t + px(3), cw, px(16),
       size=10, bold=True, color=NAVY)

rows_data = [
    ("Dijkstra",      DIJK, "0  (uninformed)", "Baseline; always optimal"),
    ("A* Manhattan",  MANH, "|Δx| + |Δy|",    "Square 4-way grids"),
    ("A* Euclidean",  EUCL, "√(Δx² + Δy²)",   "Any grid geometry"),
]
for i, (name, color, heur, best) in enumerate(rows_data):
    row_t = body_t + px(20) + i * px(26)
    if i % 2:
        rect(C1, row_t, COL_W, px(26), fill=RGBColor(0xf8, 0xf9, 0xff))
    tag = slide.shapes.add_shape(1, col_xs[0], row_t + px(4), px(110), px(17))
    tag.fill.solid(); tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tag_tf = tag.text_frame
    run = tag_tf.paragraphs[0].add_run()
    run.text = name
    run.font.size  = fsz(10)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb(heur, col_xs[1], row_t + px(5), col_ws[1], px(18),
       size=11, color=TEXT)
    tb(best, col_xs[2], row_t + px(5), col_ws[2], px(18),
       size=10, color=MUTED)

# — Conclusion & Future Work (flex:1 → ~176px) —
SEC_T += SEC_H + GAP
SEC_H = px(176)
sec_card(C1, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Conclusion & Future Work", C1, SEC_T, COL_W)
tb("Students often find informed-search concepts hard to internalize from pseudocode "
   "alone — cost trade-offs and frontier dynamics only \"click\" when made visible. "
   "This tool was built to close that gap.",
   C1 + px(10), body_t + px(6), COL_W - px(20), px(60),
   size=12, color=TEXT)
tb("The visualizer confirms heuristic choice has a measurable impact on efficiency. "
   "Learners build intuition through direct manipulation, not passive reading. "
   "Planned: hex grid support, bidirectional A*, CSV export.",
   C1 + px(10), body_t + px(72), COL_W - px(20), px(60),
   size=12, color=TEXT)

# ─────────────────────────────────────────────────────────────────────────────
# MIDDLE COLUMN
# ─────────────────────────────────────────────────────────────────────────────

# — Tool User Flow & System Architecture (flex:1.4 → ~273px) —
SEC_T = BODY_T
SEC_H = px(273)
sec_card(C2, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Tool User Flow & System Architecture", C2, SEC_T, COL_W)

DIV = C2 + px(200)
# Left: user flow
tb("User Flow".upper(), C2 + px(8), body_t + px(4), px(185), px(14),
   size=9, bold=True, color=NAVY)

flow_steps = ["START", "Configure Grid &\nPaint Terrain",
              "Choose Speed & Run", "Watch 3 Algorithms Live", "Compare Nodes & Path"]
flow_colors = [NAVY, LAV, LAV, LAV, LAV]
for i, (step, fc) in enumerate(zip(flow_steps, flow_colors)):
    sy = body_t + px(20) + i * px(42)
    shp = slide.shapes.add_shape(5, C2 + px(18), sy, px(158), px(28))  # 5 = ROUNDED_RECTANGLE
    shp.fill.solid(); shp.fill.fore_color.rgb = fc
    shp.line.color.rgb = NAVY; shp.line.width = Pt(0.5)
    run = shp.text_frame.paragraphs[0].add_run()
    run.text = step
    run.font.size = fsz(9.5)
    run.font.bold = (i == 0)
    run.font.color.rgb = WHITE if i == 0 else NAVY
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if i > 0:
        tb("↓", C2 + px(18), sy - px(14), px(158), px(14),
           size=9, color=NAVY, align=PP_ALIGN.CENTER)

# Also "Adjust & Re-run" step
tb("↕  Adjust & Re-run", C2 + px(25), body_t + px(226), px(158), px(14),
   size=8, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Divider
rect(DIV, body_t + px(2), px(1), px(SEC_H - px(33)), fill=LAV_MID)

# Right: architecture
tb("Architecture".upper(), DIV + px(8), body_t + px(4), px(242), px(14),
   size=9, bold=True, color=NAVY)

arch_items = [
    ("User Input\n(Mouse / KB)", None),
    ("Grid State\nManager",      None),
    ("Algorithm Engine ×3\n(Strategy Pattern)", None),
    ("Step Iterator / Scheduler", None),
    ("Canvas Renderer (RAF)", None),
]
for i, (label, _) in enumerate(arch_items):
    ay = body_t + px(20) + i * px(44)
    shp = slide.shapes.add_shape(1, DIV + px(10), ay, px(235), px(30))
    shp.fill.solid(); shp.fill.fore_color.rgb = LAV
    shp.line.color.rgb = LAV_MID; shp.line.width = Pt(0.5)
    run = shp.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = fsz(9.5)
    run.font.color.rgb = NAVY
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if i > 0:
        tb("↓", DIV + px(10), ay - px(14), px(235), px(14),
           size=9, color=NAVY, align=PP_ALIGN.CENTER)

# — Interactive Features (flex:1.3 → ~254px) —
SEC_T += SEC_H + GAP
SEC_H = px(254)
sec_card(C2, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Interactive Features", C2, SEC_T, COL_W)

features = [
    (RGBColor(0x37, 0x41, 0x51), "Terrain Painting",
     "Drag to paint Walls (∞), Grass (×2), or Swamp (×5) across all three panels; right-click to erase"),
    (RGBColor(0x16, 0xa3, 0x4a), "Draggable Endpoints",
     "Reposition Start / Goal at any time; re-runs automatically on next playback"),
    (RGBColor(0x63, 0x66, 0xf1), "Maze & Preset Generator",
     "DFS recursive-backtracker maze, random scatter, or barrier presets — one toolbar click"),
    (RGBColor(0x7c, 0x3a, 0xed), "Playback Controls",
     "Run / Pause / Step; four speeds from Slow (300 ms) to Max (full-batch, no delay)"),
    (AMBER,                       "f(n) / g(n) / h(n) Overlay",
     "Toggle per-cell cost values over each panel for deep algorithmic inspection"),
    (RGBColor(0xdc, 0x26, 0x26), "Live Metrics",
     "Nodes-expanded counter and path length updated live, per panel"),
]
for i, (dot_c, bold_lbl, desc) in enumerate(features):
    fy = body_t + px(5) + i * px(36)
    rect(C2 + px(8), fy + px(4), px(13), px(13), fill=dot_c)
    box = slide.shapes.add_textbox(C2 + px(26), fy, COL_W - px(32), px(34))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_lbl + " — "
    r1.font.size = fsz(11); r1.font.bold = True; r1.font.color.rgb = TEXT
    r2 = p.add_run(); r2.text = desc
    r2.font.size = fsz(11); r2.font.color.rgb = MUTED

# — Visualization Design & Color System (flex:0.9 → ~176px) —
SEC_T += SEC_H + GAP
SEC_H = px(176)
sec_card(C2, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Visualization Design & Color System", C2, SEC_T, COL_W)

tb("Purpose-built palette encodes cell states at a glance:",
   C2 + px(10), body_t + px(4), COL_W - px(20), px(16), size=11, color=TEXT)

legend_items = [
    (RGBColor(0x16, 0xa3, 0x4a), "Start (S)"),
    (RGBColor(0xdc, 0x26, 0x26), "Goal (G)"),
    (RGBColor(0x93, 0xc5, 0xfd), "Open set"),
    (RGBColor(0xfe, 0xd7, 0xaa), "Expanded"),
    (RGBColor(0x7c, 0x3a, 0xed), "Optimal path"),
    (RGBColor(0x78, 0xc3, 0x4b), "Grass (×2)"),
    (RGBColor(0x82, 0x5a, 0x28), "Swamp (×5)"),
    (RGBColor(0x26, 0x26, 0x30), "Wall (∞)"),
]
cols_per_row = 4
for i, (lc, lbl) in enumerate(legend_items):
    row, col = divmod(i, cols_per_row)
    lx = C2 + px(8) + col * px(113)
    ly = body_t + px(22) + row * px(20)
    rect(lx, ly, px(13), px(13), fill=lc)
    tb(lbl, lx + px(17), ly - px(1), px(93), px(16), size=11, color=MUTED)

tb("Algorithm panels share Start, Goal and Path colors for direct comparison. "
   "Open-set blue reveals the live frontier; expanded orange marks processed nodes. "
   "Terrain colors are visually distinct from algorithm-state colors to prevent misreading.",
   C2 + px(10), body_t + px(66), COL_W - px(20), px(54),
   size=11, color=MUTED, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# RIGHT COLUMN
# ─────────────────────────────────────────────────────────────────────────────

# — Tool Demo (flex:1.6 → ~322px) —
SEC_T = BODY_T
SEC_H = px(322)
sec_card(C3, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Tool Demo — Three Algorithms Running in Parallel", C3, SEC_T, COL_W)

# Toolbar mockup
rect(C3 + px(6), body_t + px(4), COL_W - px(12), px(22), fill=RGBColor(0x8a, 0xa4, 0xbc))
tb("Walls  Erase  Grass  Swamp  Start  Goal  Predict  Maze  Barrier  Random",
   C3 + px(10), body_t + px(7), COL_W - px(20), px(14),
   size=7.5, color=WHITE)

# Control bar
rect(C3 + px(6), body_t + px(26), COL_W - px(12), px(16),
     fill=RGBColor(0xcd, 0xd5, 0xe4))
tb("Run ▶  ‖ Pause  Step →  Reset  Clear  Slow  Med  Fast  Max  f(n) vals",
   C3 + px(10), body_t + px(28), COL_W - px(20), px(12),
   size=7, color=TEXT)

# 3 mini grid panels
panel_labels = [("Dijkstra", DIJK, "exp: 55"),
                ("A* Manhattan", MANH, "exp: 42"),
                ("A* Euclidean", EUCL, "exp: 42")]
PANEL_W  = px(131)
GRID_TOP = body_t + px(43)
GRID_H   = px(145)
CELL_N   = 10
CELL_SZ  = GRID_H / CELL_N

for pi, (pname, pcolor, pexp) in enumerate(panel_labels):
    px_start = C3 + px(6) + pi * (PANEL_W + px(5))

    # Panel header
    rect(px_start, GRID_TOP, PANEL_W, px(16), fill=pcolor)
    tb(pname, px_start + px(3), GRID_TOP + px(2), px(95), px(12),
       size=9, bold=True, color=WHITE)
    tb(pexp, px_start + px(95), GRID_TOP + px(2), px(35), px(12),
       size=8, color=WHITE, align=PP_ALIGN.RIGHT)

    # Grid cells
    for row in range(CELL_N):
        for col in range(CELL_N):
            cx_g = px_start + col * CELL_SZ
            cy_g = GRID_TOP + px(16) + row * CELL_SZ
            # cell color logic
            fc = CARD
            if col == 4 and 2 <= row <= 5:
                fc = RGBColor(0x26, 0x26, 0x30)   # wall
            elif col == 0 and row == 8:
                fc = RGBColor(0x16, 0xa3, 0x4a)   # start
            elif col == 8 and row == 1:
                fc = RGBColor(0xdc, 0x26, 0x26)   # goal
            elif col == 0 and 1 <= row <= 8:
                fc = PATH_C                        # path
            elif row == 1 and 1 <= col <= 8:
                fc = PATH_C
            elif col in (1, 2) and 3 <= row <= 4:
                fc = RGBColor(0x78, 0xc3, 0x4b)   # grass
            elif col in (5,) and 6 <= row <= 7:
                fc = RGBColor(0x82, 0x5a, 0x28)   # swamp
            elif pi == 0 and row <= 1 and col <= 9:
                fc = RGBColor(0xfe, 0xd7, 0xaa)   # dijkstra extra expanded
            elif pi > 0 and row == 0 and col <= 4:
                fc = RGBColor(0xfe, 0xd7, 0xaa)
            shp = slide.shapes.add_shape(1, cx_g + Pt(0.5), cy_g + Pt(0.5),
                                         CELL_SZ - Pt(1), CELL_SZ - Pt(1))
            shp.fill.solid(); shp.fill.fore_color.rgb = fc
            shp.line.color.rgb = RGBColor(0xe0, 0xe4, 0xf0)
            shp.line.width = Pt(0.3)

# Caption
tb("Dijkstra (left) expands more nodes than A* Manhattan (centre) and A* Euclidean (right);\nall three find the same optimal path (violet).",
   C3 + px(6), GRID_TOP + px(163), COL_W - px(12), px(28),
   size=9.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# — Performance Results & Benchmark (flex:1.4 → ~282px) —
SEC_T += SEC_H + GAP
SEC_H = px(282)
sec_card(C3, SEC_T, COL_W, SEC_H)
body_t = sec_hdr("Performance Results & Benchmark", C3, SEC_T, COL_W)

tb("Benchmark on a 10×10 square grid (seed = 42). All three algorithms return the "
   "same optimal path of length 37. Node expansion is the key efficiency differentiator:",
   C3 + px(10), body_t + px(4), COL_W - px(20), px(36),
   size=11.5, color=TEXT)

# Bar chart (manual shapes)
CHART_T = body_t + px(44)
BAR_MAX_H = px(36)
CHART_B = CHART_T + BAR_MAX_H

bars = [("Dijkstra", DIJK, 55, 55),
        ("A* Manhattan", MANH, 42, 55),
        ("A* Euclidean", EUCL, 42, 55)]
BAR_W = px(100)
BAR_GAP = px(56)

for i, (bname, bc, val, mx) in enumerate(bars):
    bh = BAR_MAX_H * val / mx
    bx = C3 + px(34) + i * (BAR_W + BAR_GAP)
    by = CHART_B - bh
    rect(bx, by, BAR_W, bh, fill=bc)
    # value label above bar
    tb(str(val), bx, by - px(16), BAR_W, px(14),
       size=12, bold=True, color=bc, align=PP_ALIGN.CENTER)

# Baseline
rect(C3 + px(20), CHART_B, COL_W - px(40), Pt(1), fill=LAV_MID)

# Bar labels
for i, (bname, bc, _, _) in enumerate(bars):
    bx = C3 + px(34) + i * (BAR_W + BAR_GAP)
    tb(bname, bx - px(10), CHART_B + px(4), BAR_W + px(20), px(14),
       size=10, bold=True, color=bc, align=PP_ALIGN.CENTER)

tb("Nodes Expanded  (lower = more efficient)",
   C3 + px(10), CHART_B + px(20), COL_W - px(20), px(12),
   size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# Summary table
TBL_T = CHART_B + px(36)
rect(C3 + px(6), TBL_T, COL_W - px(12), px(20), fill=LAV)
for j, (hdr_txt, hx) in enumerate([
        ("Algorithm",     C3 + px(10)),
        ("Path Length",   C3 + px(175)),
        ("Nodes Expanded",C3 + px(290))]):
    tb(hdr_txt, hx, TBL_T + px(3), px(130), px(14),
       size=10.5, bold=True, color=NAVY)

for i, (rname, rc, pl, ne, pct) in enumerate([
        ("Dijkstra",     DIJK, "37", "55",         ""),
        ("A* Manhattan", MANH, "37", "42  ↓ 24%",  ""),
        ("A* Euclidean", EUCL, "37", "42  ↓ 24%",  ""),
]):
    ry = TBL_T + px(20) + i * px(20)
    if i % 2 == 1:
        rect(C3 + px(6), ry, COL_W - px(12), px(20),
             fill=RGBColor(0xf9, 0xff, 0xfb))
    tb(rname, C3 + px(10), ry + px(3), px(160), px(14),
       size=11, bold=True, color=rc)
    tb(pl, C3 + px(175), ry + px(3), px(110), px(14),
       size=11, color=MUTED, align=PP_ALIGN.CENTER)
    tb(ne, C3 + px(290), ry + px(3), px(150), px(14),
       size=11, bold=(i > 0), color=rc, align=PP_ALIGN.CENTER)

# Key Insight callout
KI_T = TBL_T + px(80)
rect(C3 + px(6), KI_T, COL_W - px(12), px(90),
     fill=AMBERL, line_color=AMBER, line_w_pt=1)
tb("Key Insight".upper(), C3 + px(14), KI_T + px(6), COL_W - px(28), px(12),
   size=9.5, bold=True, color=AMBER)
tb("A* reduces node expansions by ~24% compared to Dijkstra while always "
   "guaranteeing the same optimal path — a direct, measurable benefit of "
   "informed search with an admissible heuristic.",
   C3 + px(14), KI_T + px(22), COL_W - px(28), px(62),
   size=12, bold=True, color=TEXT)

# ═════════════════════════════════════════════════════════════════════════════
# FOOTER
# ═════════════════════════════════════════════════════════════════════════════
FOOT_T = px(893)
rect(0, FOOT_T, px(1440), px(67), fill=NAVY)
tb("References",
   px(20), FOOT_T + px(6), px(200), px(16),
   size=9.5, bold=True, color=WHITE)
tb("1.  Hart, P. E., Nilsson, N. J., & Raphael, B. (1968). A formal basis for the heuristic determination of minimum cost paths. "
   "IEEE Transactions on Systems Science and Cybernetics, 4(2), 100–107.",
   px(20), FOOT_T + px(24), px(1380), px(16),
   size=9, color=RGBColor(0xc5, 0xca, 0xe9))
tb("2.  Russell, S., & Norvig, P. (2020). Artificial Intelligence: A Modern Approach (4th ed.). Pearson. Ch. 3: Solving Problems by Searching.",
   px(20), FOOT_T + px(42), px(1380), px(16),
   size=9, color=RGBColor(0xc5, 0xca, 0xe9))

# ── Save ─────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "A_Star_Poster_Editable.pptx")
prs.save(out)
print(f"Saved: {out}")

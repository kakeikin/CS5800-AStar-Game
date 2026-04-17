"""
Generate A* Pathfinding Visualizer presentation PPTX
Run: python3 generate_pptx.py
Output: web/A_Star_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1a, 0x23, 0x7e)
NAVY2   = RGBColor(0x28, 0x35, 0x93)
LAV     = RGBColor(0xe8, 0xea, 0xf6)
WHITE   = RGBColor(0xff, 0xff, 0xff)
TEXT    = RGBColor(0x1c, 0x1c, 0x2e)
MUTED   = RGBColor(0x4a, 0x4a, 0x6a)
DIJK    = RGBColor(0x02, 0x84, 0xc7)
MANH    = RGBColor(0x05, 0x96, 0x69)
EUCL    = RGBColor(0xdc, 0x68, 0x03)
AMBER   = RGBColor(0xd9, 0x77, 0x06)
AMBERL  = RGBColor(0xff, 0xf7, 0xe6)
GRAY_BG = RGBColor(0xf0, 0xf2, 0xfc)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def header_bar(slide, title_text, subtitle=None):
    """Dark navy top bar."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=NAVY)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.12), Inches(10), Inches(0.7),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.78), Inches(10), Inches(0.35),
                 size=14, color=RGBColor(0xc5, 0xca, 0xe9))
    add_text(slide, "Northeastern University",
             Inches(10.6), Inches(0.35), Inches(2.5), Inches(0.55),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

def section_card(slide, title, l, t, w, h, title_color=NAVY):
    """Lavender-header section card."""
    add_rect(slide, l, t, w, h, fill=WHITE,
             line=RGBColor(0xd0, 0xd4, 0xee), line_w=Pt(0.5))
    add_rect(slide, l, t, w, Inches(0.34), fill=LAV)
    add_text(slide, title.upper(),
             l + Inches(0.1), t + Inches(0.04), w - Inches(0.2), Inches(0.28),
             size=10, bold=True, color=title_color)
    return (l + Inches(0.15), t + Inches(0.38))  # body origin

def bullet(slide, text, l, t, w, size=13, indent=False, bold_prefix=None):
    tb = slide.shapes.add_textbox(l + (Inches(0.2) if indent else 0), t, w, Inches(0.38))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix + " — "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = TEXT
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(size)
        r2.font.color.rgb = MUTED
    else:
        run = p.add_run()
        run.text = ("• " if not indent else "  – ") + text
        run.font.size = Pt(size)
        run.font.color.rgb = MUTED
    return t + Inches(0.37)

# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
add_rect(sl, 0, 0, SLIDE_W, Inches(3.8), fill=NAVY)
add_rect(sl, 0, Inches(3.8), SLIDE_W, Inches(0.06), fill=AMBER)

add_text(sl, "A* Pathfinding Visualizer",
         Inches(0.8), Inches(0.5), Inches(11), Inches(1.5),
         size=52, bold=True, color=WHITE)
add_text(sl, "Interactive side-by-side comparison of Dijkstra, A* Manhattan, and A* Euclidean\non weighted grids with terrain, step-through replay, and f/g/h overlays",
         Inches(0.8), Inches(2.05), Inches(9.5), Inches(1.0),
         size=16, color=RGBColor(0xc5, 0xca, 0xe9))
add_text(sl, "Jiaxin Jia  ·  Xiaoyuan Lu  ·  Xinyuan Fan",
         Inches(0.8), Inches(4.15), Inches(8), Inches(0.45),
         size=15, bold=True, color=TEXT)
add_text(sl, "Khoury College of Computer Sciences  ·  Northeastern University",
         Inches(0.8), Inches(4.6), Inches(8), Inches(0.4),
         size=13, italic=True, color=MUTED)
add_text(sl, "CS5800 — Spring 2026",
         Inches(0.8), Inches(5.05), Inches(8), Inches(0.4),
         size=13, color=MUTED)

add_rect(sl, Inches(10.0), Inches(4.0), Inches(3.0), Inches(3.1), fill=WHITE,
         line=RGBColor(0xd0, 0xd4, 0xee), line_w=Pt(0.5))
for i, (label, color) in enumerate([
        ("Dijkstra", DIJK), ("A* Manhattan", MANH), ("A* Euclidean", EUCL)]):
    add_rect(sl, Inches(10.2), Inches(4.25) + i*Inches(0.82),
             Inches(0.18), Inches(0.18), fill=color)
    add_text(sl, label,
             Inches(10.5), Inches(4.2) + i*Inches(0.82),
             Inches(2.3), Inches(0.3), size=13, bold=True, color=color)

# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — Background & Motivation
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "Background & Motivation",
           "Why build a side-by-side pathfinding visualizer?")

cx1, cy1 = section_card(sl, "The Problem",
                         Inches(0.3), Inches(1.3), Inches(6.1), Inches(2.7))
points_l = [
    "Algorithm trade-offs are hard to internalize from pseudocode alone",
    "The contrast between uninformed and informed search only becomes intuitive when you can watch both exploration frontiers side by side",
    "Most tools show one algorithm at a time on uniform-cost grids — direct comparison is impossible",
]
y = cy1
for p in points_l:
    y = bullet(sl, p, cx1, y, Inches(5.7), size=13)

cx2, cy2 = section_card(sl, "Our Solution",
                          Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.7))
points_r = [
    "Run Dijkstra, A* Manhattan, and A* Euclidean simultaneously on shared weighted grids",
    "Step-by-step replay with per-cell f/g/h overlays reveal exactly why each algorithm expands each node",
    "Student feedback confirmed: watching A*'s frontier stay focused while Dijkstra fans out 'clicks' in a way no diagram can replicate",
]
y = cy2
for p in points_r:
    y = bullet(sl, p, cx2, y, Inches(5.9), size=13)

add_rect(sl, Inches(0.3), Inches(4.2), Inches(12.73), Inches(1.2),
         fill=AMBERL, line=AMBER, line_w=Pt(1))
add_text(sl, "KEY MOTIVATION",
         Inches(0.5), Inches(4.28), Inches(3), Inches(0.28),
         size=9, bold=True, color=AMBER)
add_text(sl, "Learning algorithms by watching them — not reading about them.",
         Inches(0.5), Inches(4.56), Inches(12.2), Inches(0.6),
         size=16, bold=True, color=TEXT)

# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — Algorithm Overview
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "Algorithm Overview",
           "Three search strategies — same grid, same start/goal")

cols = [
    ("Dijkstra", DIJK, "0  (uninformed)",
     "Baseline; explores uniformly in all directions regardless of goal location",
     "Blind search — no knowledge of goal position"),
    ("A* Manhattan", MANH, "|Δx| + |Δy|",
     "Optimized for 4-directional square grids; guides search toward goal using grid distance",
     "Underestimates on diagonal or free-movement grids"),
    ("A* Euclidean", EUCL, "√(Δx² + Δy²)",
     "Straight-line distance heuristic; works for any movement geometry",
     "Slightly more exploration than Manhattan on pure grid movement"),
]
for i, (name, color, h, desc, note) in enumerate(cols):
    lx = Inches(0.3) + i * Inches(4.34)
    add_rect(sl, lx, Inches(1.3), Inches(4.1), Inches(5.7), fill=WHITE,
             line=RGBColor(0xd0, 0xd4, 0xee), line_w=Pt(0.5))
    add_rect(sl, lx, Inches(1.3), Inches(4.1), Inches(0.65), fill=color)
    add_text(sl, name, lx + Inches(0.15), Inches(1.35),
             Inches(3.8), Inches(0.55), size=20, bold=True, color=WHITE)
    add_text(sl, "Heuristic h(n)", lx + Inches(0.15), Inches(2.1),
             Inches(3.8), Inches(0.28), size=10, bold=True, color=MUTED)
    add_text(sl, h, lx + Inches(0.15), Inches(2.38),
             Inches(3.8), Inches(0.4), size=18, bold=True, color=color)
    add_text(sl, "How it works", lx + Inches(0.15), Inches(2.88),
             Inches(3.8), Inches(0.28), size=10, bold=True, color=MUTED)
    add_text(sl, desc, lx + Inches(0.15), Inches(3.16),
             Inches(3.8), Inches(1.0), size=12, color=TEXT)
    add_rect(sl, lx + Inches(0.1), Inches(4.4), Inches(3.9), Inches(0.45), fill=LAV)
    add_text(sl, "Note: " + note, lx + Inches(0.2), Inches(4.45),
             Inches(3.7), Inches(0.38), size=11, italic=True, color=MUTED)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.28), fill=LAV)
add_text(sl, "All three use:  f(n) = g(n) + h(n)   where g(n) = cost so far, h(n) = heuristic estimate to goal.  Setting h(n) = 0 recovers Dijkstra.",
         Inches(0.5), Inches(7.12), Inches(12.2), Inches(0.24),
         size=11, color=MUTED, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — System Architecture & Features
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "System Architecture & Interactive Features")

cx, cy = section_card(sl, "System Architecture",
                       Inches(0.3), Inches(1.3), Inches(4.8), Inches(5.7))
arch_steps = [
    ("User Input", "Mouse / Keyboard events captured in real time"),
    ("Grid State Manager", "Shared terrain grid synced across all 3 panels"),
    ("Algorithm Engine ×3", "Independent instances run in lockstep"),
    ("Step Iterator / Scheduler", "Controls playback speed and step-through"),
    ("Canvas Renderer (RAF)", "requestAnimationFrame draws each frame"),
]
y = cy
for i, (title, desc) in enumerate(arch_steps):
    if i > 0:
        add_text(sl, "↓", cx, y - Inches(0.22), Inches(4.5), Inches(0.25),
                 size=14, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(sl, cx, y, Inches(4.5), Inches(0.7),
             fill=NAVY if i == 0 else LAV,
             line=RGBColor(0xc5, 0xca, 0xe9), line_w=Pt(0.5))
    add_text(sl, title, cx + Inches(0.1), y + Inches(0.04),
             Inches(4.3), Inches(0.3), size=12, bold=True,
             color=WHITE if i == 0 else NAVY)
    add_text(sl, desc, cx + Inches(0.1), y + Inches(0.33),
             Inches(4.3), Inches(0.3), size=10,
             color=WHITE if i == 0 else MUTED)
    y += Inches(0.95)

cx2, cy2 = section_card(sl, "Interactive Features",
                          Inches(5.4), Inches(1.3), Inches(7.6), Inches(5.7))
features = [
    ("Terrain Painting", "Paint Walls (∞), Grass (×2), Swamp (×5); drag across all 3 panels; right-click to erase"),
    ("Draggable Endpoints", "Reposition Start / Goal at any time; re-runs automatically"),
    ("Maze & Preset Generator", "DFS recursive-backtracker maze, random scatter, or barrier presets — one click"),
    ("Playback Controls", "Run / Pause / Step; four speeds from Slow (300ms) to Max (full-batch, no delay)"),
    ("f(n) / g(n) / h(n) Overlay", "Toggle per-cell cost values over each panel for deep algorithmic inspection"),
    ("Live Metrics", "Nodes-expanded counter and path length updated live per panel"),
]
y = cy2
for bold, desc in features:
    y = bullet(sl, desc, cx2, y, Inches(7.2), size=12, bold_prefix=bold)
    y += Inches(0.04)

# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Tool Demo
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "Tool Demo — Three Algorithms Running in Parallel",
           "Same start, same goal, same terrain — different exploration strategies")

panel_names = [("Dijkstra", DIJK), ("A* Manhattan", MANH), ("A* Euclidean", EUCL)]
for i, (name, color) in enumerate(panel_names):
    px = Inches(0.3) + i * Inches(4.34)
    add_rect(sl, px, Inches(1.3), Inches(4.1), Inches(0.45), fill=color)
    add_text(sl, name, px + Inches(0.15), Inches(1.33),
             Inches(3.8), Inches(0.38), size=16, bold=True, color=WHITE)
    add_rect(sl, px, Inches(1.75), Inches(4.1), Inches(4.2), fill=WHITE,
             line=RGBColor(0xc0, 0xc8, 0xe0), line_w=Pt(0.5))

    cell = Inches(0.35)
    for row in range(10):
        for col in range(10):
            cx2 = px + Inches(0.1) + col * cell
            cy2 = Inches(1.8) + row * cell
            fill_c = WHITE
            if col == 4 and 2 <= row <= 5:
                fill_c = RGBColor(0x26, 0x26, 0x30)
            elif col == 0 and row == 8:
                fill_c = RGBColor(0x16, 0xa3, 0x4a)
            elif col == 8 and row == 1:
                fill_c = RGBColor(0xdc, 0x26, 0x26)
            elif col == 0 and 1 <= row <= 8:
                fill_c = RGBColor(0x7c, 0x3a, 0xed)
            elif row == 1 and 0 <= col <= 8:
                fill_c = RGBColor(0x7c, 0x3a, 0xed)
            elif i == 0 and row <= 1:
                if col <= 8: fill_c = RGBColor(0xfe, 0xd7, 0xaa)
            elif row <= 1 and col <= 5:
                fill_c = RGBColor(0xfe, 0xd7, 0xaa)
            add_rect(sl, cx2, cy2, cell - Inches(0.02), cell - Inches(0.02),
                     fill=fill_c, line=RGBColor(0xe0, 0xe4, 0xf0), line_w=Pt(0.3))

legend_items = [
    ("Start (S)", RGBColor(0x16, 0xa3, 0x4a)),
    ("Goal (G)",  RGBColor(0xdc, 0x26, 0x26)),
    ("Optimal path", RGBColor(0x7c, 0x3a, 0xed)),
    ("Expanded", RGBColor(0xfe, 0xd7, 0xaa)),
    ("Wall (∞)", RGBColor(0x26, 0x26, 0x30)),
]
for i, (label, color) in enumerate(legend_items):
    lx = Inches(0.5) + i * Inches(2.55)
    add_rect(sl, lx, Inches(6.3), Inches(0.22), Inches(0.22), fill=color)
    add_text(sl, label, lx + Inches(0.28), Inches(6.27),
             Inches(2.1), Inches(0.28), size=11, color=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Performance Results
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "Performance Results & Benchmark",
           "10×10 square grid, seed = 42  —  identical start / goal / terrain")

chart_data = ChartData()
chart_data.categories = ['Dijkstra', 'A* Manhattan', 'A* Euclidean']
chart_data.add_series('Nodes Expanded', (55, 42, 42))

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.3), Inches(6.5), Inches(4.8),
    chart_data
).chart
chart.has_title = True
chart.chart_title.text_frame.text = "Nodes Expanded"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT
chart.has_legend = False

from pptx.oxml.ns import qn
from lxml import etree
ser = chart.plots[0].series[0]
pts_xml = ser._element
for idx, color in enumerate([DIJK, MANH, EUCL]):
    dp = etree.SubElement(pts_xml, qn('c:dPt'))
    idx_el = etree.SubElement(dp, qn('c:idx'))
    idx_el.set('val', str(idx))
    spPr = etree.SubElement(dp, qn('c:spPr'))
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{color[0]:02x}{color[1]:02x}{color[2]:02x}')

cx, cy = section_card(sl, "Results Summary",
                       Inches(7.2), Inches(1.3), Inches(5.8), Inches(3.2))
add_rect(sl, Inches(7.2), cy - Inches(0.05), Inches(5.8), Inches(0.4), fill=LAV)
for j, (hdr, lx) in enumerate([
        ("Algorithm", Inches(7.3)),
        ("Path Length", Inches(9.5)),
        ("Nodes Expanded", Inches(11.0))]):
    add_text(sl, hdr, lx, cy - Inches(0.02), Inches(1.8), Inches(0.36),
             size=11, bold=True, color=NAVY)

rows = [
    ("Dijkstra",     DIJK, "37", "55"),
    ("A* Manhattan", MANH, "37", "42  ↓ 24%"),
    ("A* Euclidean", EUCL, "37", "42  ↓ 24%"),
]
y = cy + Inches(0.42)
for i, (name, color, pl, ne) in enumerate(rows):
    bg = RGBColor(0xf9, 0xff, 0xfb) if i % 2 else WHITE
    add_rect(sl, Inches(7.2), y, Inches(5.8), Inches(0.45), fill=bg,
             line=RGBColor(0xf0, 0xf0, 0xf8), line_w=Pt(0.3))
    add_text(sl, name, Inches(7.3), y + Inches(0.05),
             Inches(2.1), Inches(0.36), size=12, bold=True, color=color)
    add_text(sl, pl,   Inches(9.5), y + Inches(0.05),
             Inches(1.4), Inches(0.36), size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, ne,   Inches(11.0), y + Inches(0.05),
             Inches(1.7), Inches(0.36), size=12, bold=(i>0), color=color, align=PP_ALIGN.CENTER)
    y += Inches(0.47)

add_rect(sl, Inches(7.2), Inches(4.7), Inches(5.8), Inches(2.1),
         fill=AMBERL, line=AMBER, line_w=Pt(1.5))
add_text(sl, "KEY INSIGHT",
         Inches(7.4), Inches(4.78), Inches(5.4), Inches(0.3),
         size=10, bold=True, color=AMBER)
add_text(sl, "A* reduces node expansions by ~24%\ncompared to Dijkstra while always\nguaranteeing the same optimal path.",
         Inches(7.4), Inches(5.1), Inches(5.4), Inches(1.4),
         size=16, bold=True, color=TEXT)

# ════════════════════════════════════════════════════════════════════════════
# Slide 7 — Conclusion & Future Work
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_BG)
header_bar(sl, "Conclusion & Future Work")

cx, cy = section_card(sl, "What We Built & Learned",
                       Inches(0.3), Inches(1.3), Inches(6.1), Inches(5.7))
conc_points = [
    "Built a browser-based tool that runs three pathfinding algorithms simultaneously on shared weighted grids",
    "Confirmed that heuristic choice has a measurable ~24% efficiency impact on node expansions without sacrificing path quality",
    "Demonstrated that direct manipulation and live visualization are more effective for building intuition than static diagrams",
    "Terrain weighting (Grass ×2, Swamp ×5) adds real-world complexity that textbook uniform grids omit",
]
y = cy
for p in conc_points:
    y = bullet(sl, p, cx, y, Inches(5.7), size=13)
    y += Inches(0.06)

cx2, cy2 = section_card(sl, "Future Work",
                          Inches(6.7), Inches(1.3), Inches(6.3), Inches(5.7))
future = [
    ("Hex Grid Support", "Extend to hexagonal grids for 6-directional movement"),
    ("Bidirectional A*", "Run search from both Start and Goal simultaneously to further reduce expansions"),
    ("CSV Export", "Export benchmark data (nodes expanded, path cost, time) for offline analysis"),
    ("More Algorithms", "Add Greedy Best-First, JPS (Jump Point Search), and weighted A* variants"),
    ("Shareable Configs", "URL-encoded grid state so users can share specific puzzle scenarios"),
]
y = cy2
for bold, desc in future:
    y = bullet(sl, desc, cx2, y, Inches(6.0), size=12, bold_prefix=bold)
    y += Inches(0.08)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "A_Star_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
